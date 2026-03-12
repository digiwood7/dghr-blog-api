"""
파일 텍스트 추출 서비스
PDF, DOCX, PPTX, XLSX, TXT에서 텍스트를 추출하여 구조화된 JSON 반환
"""
import io
import re
from pathlib import Path

# PDF
try:
    from PyPDF2 import PdfReader
    HAS_PYPDF2 = True
except ImportError:
    HAS_PYPDF2 = False

# DOCX
try:
    from docx import Document as DocxDocument
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

# PPTX
try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

# XLSX
try:
    from openpyxl import load_workbook
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


def extract_text_from_pdf(content: bytes, max_chars: int = 10000) -> dict:
    """PDF에서 텍스트 추출"""
    if not HAS_PYPDF2:
        return {"error": "PyPDF2 미설치", "text": "", "pages": 0}

    try:
        reader = PdfReader(io.BytesIO(content))
        pages = []
        total_text = ""
        for i, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            text = re.sub(r'\s+', ' ', text).strip()
            if text:
                pages.append({"page": i + 1, "text": text[:2000]})
                total_text += text + "\n"
            if len(total_text) >= max_chars:
                break

        return {
            "text": total_text[:max_chars],
            "pages": len(reader.pages),
            "extracted_pages": len(pages),
        }
    except Exception as e:
        return {"error": str(e), "text": "", "pages": 0}


def extract_text_from_docx(content: bytes, max_chars: int = 10000) -> dict:
    """DOCX에서 텍스트 추출"""
    if not HAS_DOCX:
        return {"error": "python-docx 미설치", "text": ""}

    try:
        doc = DocxDocument(io.BytesIO(content))
        paragraphs = []
        total_text = ""
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                paragraphs.append(text)
                total_text += text + "\n"
            if len(total_text) >= max_chars:
                break

        # 테이블 내용도 추출
        for table in doc.tables:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    total_text += row_text + "\n"
                if len(total_text) >= max_chars:
                    break

        return {
            "text": total_text[:max_chars],
            "paragraphs": len(paragraphs),
        }
    except Exception as e:
        return {"error": str(e), "text": ""}


def extract_text_from_pptx(content: bytes, max_chars: int = 10000) -> dict:
    """PPTX에서 텍스트 추출 (슬라이드별 구조화)"""
    if not HAS_PPTX:
        return {"error": "python-pptx 미설치", "text": "", "slides": 0}

    try:
        prs = Presentation(io.BytesIO(content))
        slides_data = []
        total_text = ""
        for i, slide in enumerate(prs.slides):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        text = paragraph.text.strip()
                        if text:
                            slide_texts.append(text)
                if hasattr(shape, "table"):
                    for row in shape.table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                        if row_text:
                            slide_texts.append(row_text)

            if slide_texts:
                slide_content = "\n".join(slide_texts)
                slides_data.append({"slide": i + 1, "text": slide_content[:1000]})
                total_text += f"[슬라이드 {i+1}]\n{slide_content}\n\n"
            if len(total_text) >= max_chars:
                break

        return {
            "text": total_text[:max_chars],
            "slides": len(prs.slides),
            "extracted_slides": len(slides_data),
        }
    except Exception as e:
        return {"error": str(e), "text": "", "slides": 0}


def extract_text_from_xlsx(content: bytes, max_chars: int = 10000) -> dict:
    """XLSX에서 텍스트 추출"""
    if not HAS_OPENPYXL:
        return {"error": "openpyxl 미설치", "text": ""}

    try:
        wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
        total_text = ""
        sheets_data = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            sheet_text = f"[시트: {sheet_name}]\n"
            row_count = 0
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) if c is not None else "" for c in row]
                row_str = " | ".join(c for c in cells if c)
                if row_str:
                    sheet_text += row_str + "\n"
                    row_count += 1
                if len(total_text) + len(sheet_text) >= max_chars:
                    break
            sheets_data.append({"sheet": sheet_name, "rows": row_count})
            total_text += sheet_text + "\n"
            if len(total_text) >= max_chars:
                break
        wb.close()

        return {
            "text": total_text[:max_chars],
            "sheets": len(wb.sheetnames),
        }
    except Exception as e:
        return {"error": str(e), "text": ""}


def extract_text_from_txt(content: bytes, max_chars: int = 10000) -> dict:
    """TXT 파일에서 텍스트 추출"""
    try:
        for encoding in ["utf-8", "euc-kr", "cp949", "latin-1"]:
            try:
                text = content.decode(encoding)
                return {"text": text[:max_chars]}
            except (UnicodeDecodeError, LookupError):
                continue
        return {"error": "인코딩 감지 실패", "text": ""}
    except Exception as e:
        return {"error": str(e), "text": ""}


def extract_file_content(content: bytes, filename: str, max_chars: int = 10000) -> dict:
    """파일 확장자에 따라 적절한 추출기 호출"""
    ext = Path(filename).suffix.lower()

    extractors = {
        ".pdf": extract_text_from_pdf,
        ".docx": extract_text_from_docx,
        ".doc": extract_text_from_docx,  # doc은 지원 한계 있음
        ".pptx": extract_text_from_pptx,
        ".ppt": extract_text_from_pptx,
        ".xlsx": extract_text_from_xlsx,
        ".xls": extract_text_from_xlsx,
        ".txt": extract_text_from_txt,
    }

    extractor = extractors.get(ext)
    if not extractor:
        return {
            "filename": filename,
            "file_type": ext,
            "text": "",
            "error": f"텍스트 추출 미지원 형식: {ext}",
        }

    result = extractor(content, max_chars)
    result["filename"] = filename
    result["file_type"] = ext
    return result
